"""
Native loaders for non-PDF formats → the SAME typed-block schema the PDF parser
emits. This lets Office/text documents flow through the identical downstream
pipeline: doc-type routing → structure-aware chunking → hybrid retrieval → rerank.

Block schema (matches rag.parsers / block_chunker):
  {"type": "heading"|"paragraph"|"table"|"list", "text"|"md": ..., "level"?: int, "page": int}

Returns None for unsupported extensions so the caller can fall back to LlamaParse.
"""
import os
import csv as _csv
from typing import Any, Dict, List, Optional

import structlog

logger = structlog.get_logger(__name__)

SUPPORTED_EXTS = {".docx", ".pptx", ".xlsx", ".xls", ".csv", ".txt", ".md", ".markdown"}
MAX_TABLE_ROWS = int(os.getenv("NATIVE_MAX_TABLE_ROWS", "200"))


def blocks_from_file(path: str, filename: str) -> Optional[List[Dict[str, Any]]]:
    ext = os.path.splitext(filename)[1].lower()
    try:
        if ext == ".docx":
            return _docx_blocks(path)
        if ext == ".pptx":
            return _pptx_blocks(path)
        if ext in (".xlsx", ".xls"):
            return _xlsx_blocks(path)
        if ext == ".csv":
            return _csv_blocks(path)
        if ext in (".txt", ".md", ".markdown"):
            return _text_blocks(path)
    except Exception as e:
        logger.warning("native block extraction failed", ext=ext, error=str(e))
        return None
    return None


def _rows_to_md(rows: List[List[str]]) -> str:
    rows = [[("" if c is None else str(c)).replace("\n", " ").strip() for c in r] for r in rows if r]
    rows = [r for r in rows if any(c for c in r)]
    if not rows:
        return ""
    ncol = max(len(r) for r in rows)
    rows = [r + [""] * (ncol - len(r)) for r in rows]
    header, body = rows[0], rows[1:]
    out = ["| " + " | ".join(header) + " |", "| " + " | ".join(["---"] * ncol) + " |"]
    for r in body[:MAX_TABLE_ROWS]:
        out.append("| " + " | ".join(r) + " |")
    if len(body) > MAX_TABLE_ROWS:
        out.append(f"| … {len(body) - MAX_TABLE_ROWS} more rows … |")
    return "\n".join(out)


def _docx_blocks(path: str) -> List[Dict[str, Any]]:
    from docx import Document as Docx
    from docx.oxml.ns import qn
    doc = Docx(path)
    tables = iter(doc.tables)
    blocks: List[Dict[str, Any]] = []
    # iterate body children in document order (paragraphs AND tables interleaved)
    for child in doc.element.body.iterchildren():
        if child.tag == qn("w:p"):
            from docx.text.paragraph import Paragraph
            para = Paragraph(child, doc)
            text = para.text.strip()
            if not text:
                continue
            style = (para.style.name or "").lower() if para.style else ""
            if style.startswith("heading") or style == "title":
                digits = "".join(ch for ch in style if ch.isdigit())
                lvl = int(digits) if digits else 1
                blocks.append({"type": "heading", "level": min(max(lvl, 1), 6), "text": text, "page": 1})
            else:
                blocks.append({"type": "paragraph", "text": text, "page": 1})
        elif child.tag == qn("w:tbl"):
            try:
                tbl = next(tables)
                md = _rows_to_md([[cell.text for cell in row.cells] for row in tbl.rows])
                if md:
                    blocks.append({"type": "table", "md": md, "page": 1})
            except StopIteration:
                pass
    return blocks


def _pptx_blocks(path: str) -> List[Dict[str, Any]]:
    from pptx import Presentation
    prs = Presentation(path)
    blocks: List[Dict[str, Any]] = []
    for i, slide in enumerate(prs.slides, start=1):
        title = ""
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            title = slide.shapes.title.text.strip()
        blocks.append({"type": "heading", "level": 1, "text": title or f"Slide {i}", "page": i})
        for shape in slide.shapes:
            if shape == slide.shapes.title or not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if text:
                blocks.append({"type": "paragraph", "text": text, "page": i})
    return blocks


def _xlsx_blocks(path: str) -> List[Dict[str, Any]]:
    from openpyxl import load_workbook
    wb = load_workbook(path, read_only=True, data_only=True)
    blocks: List[Dict[str, Any]] = []
    for si, ws in enumerate(wb.worksheets, start=1):
        blocks.append({"type": "heading", "level": 1, "text": ws.title, "page": si})
        rows = [[c for c in row] for row in ws.iter_rows(values_only=True)]
        md = _rows_to_md(rows)
        if md:
            blocks.append({"type": "table", "md": md, "page": si})
    wb.close()
    return blocks


def _csv_blocks(path: str) -> List[Dict[str, Any]]:
    with open(path, "r", encoding="utf-8", errors="replace", newline="") as f:
        rows = list(_csv.reader(f))
    md = _rows_to_md(rows)
    return [{"type": "table", "md": md, "page": 1}] if md else []


def _text_blocks(path: str) -> List[Dict[str, Any]]:
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        text = f.read()
    blocks: List[Dict[str, Any]] = []
    para: List[str] = []

    def flush_para():
        if para:
            joined = " ".join(para).strip()
            if joined:
                blocks.append({"type": "paragraph", "text": joined, "page": 1})
            para.clear()

    for line in text.splitlines():
        s = line.strip()
        if s.startswith("#"):                    # markdown heading
            flush_para()
            lvl = len(s) - len(s.lstrip("#"))
            blocks.append({"type": "heading", "level": min(max(lvl, 1), 6),
                           "text": s.lstrip("#").strip(), "page": 1})
        elif not s:                              # blank line = paragraph break
            flush_para()
        else:
            para.append(s)
    flush_para()
    return blocks
